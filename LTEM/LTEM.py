# 将文件夹内部的 mrc 和 dm4 文件做成 PPT — ttkbootstrap 重写版
import os
import time
import threading
from pathlib import Path
import numpy as np
import matplotlib
matplotlib.use('TkAgg')
from matplotlib.backends.backend_tkagg import FigureCanvasTkAgg
from matplotlib.figure import Figure
import mrcfile
import hyperspy.api as hs
import imageio.v2 as imageio
import cv2
from PIL import Image
import tkinter as tk
from tkinter import filedialog, messagebox
import ttkbootstrap as ttk
from pptx import Presentation
from pptx.util import Cm

BASE_DIR = Path(__file__).resolve().parent
GIF_PATH = BASE_DIR / "neco4.gif"
GIF_DISPLAY_SCALE = 0.80


# ============================================================
# 数据加载
# ============================================================
def load_data(filepath):
    """读取 mrc/dm4，保留原始维度"""
    ext = filepath.lower().split('.')[-1]
    if ext == 'mrc':
        with mrcfile.open(filepath) as mrc:
            data = mrc.data.copy()
    elif ext == 'dm4':
        img = hs.load(filepath)
        data = img.data
    else:
        raise ValueError(f"不支持的文件格式: {ext}")
    if np.iscomplexobj(data):
        data = np.abs(data)
    return data


def load_preview_frame(filepath):
    """读取文件的第一帧用于预览 (2D)"""
    data = load_data(filepath)
    if data.ndim == 3:
        data = data[0]
    data = np.rot90(data, 2)
    data = np.fliplr(data)
    return data


def apply_bc(data, low_percentile, high_percentile, brightness, contrast):
    """Fiji 风格的亮度/对比度调整"""
    lo = np.percentile(data, max(0.001, low_percentile))
    hi = np.percentile(data, min(99.999, high_percentile))
    if hi <= lo:
        hi = lo + 1e-6
    result = np.clip(data, lo, hi)
    result = (result - lo) / (hi - lo)
    result = result + brightness
    if contrast > 0:
        result = (result - 0.5) * contrast + 0.5
    result = np.clip(result, 0, 1)
    return result, lo, hi


def make_video(data_3d, video_path, low_pct, high_pct, brightness, contrast,
               filename="", fps=5, scale_factor=1, show_info=True):
    """将多层数据做成 mp4 视频，应用 B&C 调整。
    show_info: 是否在画面左上角叠加文件名 + 帧数信息。
    """
    frames = []
    n_frames = data_3d.shape[0]
    for i in range(0, n_frames):
        frame = data_3d[i]
        if np.iscomplexobj(frame):
            frame = np.abs(frame)
        frame = np.rot90(frame, 2)
        frame = np.fliplr(frame)
        frame, _, _ = apply_bc(frame, low_pct, high_pct, brightness, contrast)
        if scale_factor != 1:
            frame = cv2.resize(frame, None, fx=scale_factor, fy=scale_factor,
                               interpolation=cv2.INTER_CUBIC)
        frame_uint8 = (np.clip(frame, 0, 1) * 255).astype(np.uint8)
        frame_rgb = np.stack([frame_uint8] * 3, axis=-1)
        if show_info:
            text = f"{filename} | frame {i + 1}/{n_frames}"
            cv2.putText(frame_rgb, text, (20, 40), cv2.FONT_HERSHEY_SIMPLEX,
                        1, (255, 255, 255), 2, cv2.LINE_AA)
        frames.append(frame_rgb)
    imageio.mimsave(video_path, frames, fps=fps)
    return video_path


# ============================================================
# PPT 生成（独立函数，便于后台线程调用）
# ============================================================
def generate_ppt(folder, files, mode, resolution, per_file_params, default_params,
                 progress_cb=None, show_info=True):
    """根据当前参数生成 PPT（在后台线程中运行）。
    show_info: 多层数据生成的 mp4 视频是否叠加文件名 + 帧数信息。
    """
    folder = str(folder)
    out_scale = resolution / 100

    if len(files) >= 2:
        second_file_path = os.path.join(folder, files[1])
    else:
        second_file_path = folder

    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    cols = 6
    cell_h = Cm(5.5)
    cell_w = Cm(5.6)
    x_gap = Cm(0.01)
    left_margin = Cm(0.1)
    top_margin = Cm(1.5)

    total = len(files)
    for idx, file in enumerate(files):
        if idx % 12 == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            textbox_top = slide.shapes.add_textbox(Cm(1), Cm(0.2), Cm(30), Cm(0.6))
            textbox_top.text = second_file_path

        filepath = os.path.join(folder, file)
        name_no_ext = os.path.splitext(file)[0]

        raw_data = load_data(filepath)
        is_multi = (raw_data.ndim == 3 and raw_data.shape[0] > 1)

        if mode == 'individual' and file in per_file_params:
            p = per_file_params[file]
            _lp, _hp, _br, _ct = p['low_pct'], p['high_pct'], p['brightness'], p['contrast']
        else:
            _lp, _hp, _br, _ct = (default_params['low_pct'], default_params['high_pct'],
                                   default_params['brightness'], default_params['contrast'])

        pos = idx % 12
        row = pos // cols
        col = pos % cols
        x = left_margin + col * (cell_w + x_gap)
        y = top_margin if row == 0 else top_margin + Cm(9)
        title_y = y if col % 2 == 0 else y + Cm(1.2)

        textbox = slide.shapes.add_textbox(x, title_y, cell_w, Cm(0.5))
        textbox.text = name_no_ext

        if is_multi:
            video_path = os.path.join(folder, f"{name_no_ext}.mp4")
            make_video(raw_data, video_path, _lp, _hp, _br, _ct,
                       filename=file, scale_factor=out_scale, show_info=show_info)
            slide.shapes.add_movie(video_path, x, y + Cm(2),
                                   width=cell_w, height=cell_h, poster_frame_image=None)
        else:
            if raw_data.ndim == 3:
                raw_data = raw_data[0]
            data_2d = np.rot90(raw_data, 2)
            data_2d = np.fliplr(data_2d)
            if np.iscomplexobj(data_2d):
                data_2d = np.abs(data_2d)
            data_2d, _, _ = apply_bc(data_2d, _lp, _hp, _br, _ct)

            img_path = os.path.join(folder, f"{name_no_ext}.png")
            img_uint8 = (data_2d * 255).astype(np.uint8)
            img = Image.fromarray(img_uint8)
            if out_scale != 1.0:
                new_w = int(img.width * out_scale)
                new_h = int(img.height * out_scale)
                img = img.resize((new_w, new_h), Image.Resampling.LANCZOS)
            img.save(img_path)
            slide.shapes.add_picture(img_path, x, y + Cm(2), height=cell_h)

        if progress_cb is not None:
            progress_cb(idx + 1, total, file)

    out_path = os.path.join(folder, "mrc_dm4_output_bc.pptx")
    prs.save(out_path)
    return out_path


# ============================================================
# ttkbootstrap GUI
# ============================================================
class App(ttk.Window):
    def __init__(self):
        super().__init__(themename="flatly")
        self.title("MRC/DM4 → PPT — 亮度/对比度调整")
        self.geometry("1000x900")
        self.resizable(True, True)

        # 状态
        self.mode = 'fixed'
        self.resolution = 50
        self.show_info = True
        self.folder = None
        self.files = []
        self.current_file_idx = 0
        self.per_file_params = {}
        self.preview_data = None
        self.sorted_flat = None
        self.last_update_time = 0

        self.buttons = []
        self.main_frame = None
        self.run_frame = None
        self.run_gif_label = None
        self.run_gif_frames = []
        self.gif_running = False
        self.gif_after_id = None
        self._preloaded_frames = self.load_gif(GIF_PATH)

        self.build_ui()

    # ---------- UI 搭建 ----------
    def build_ui(self):
        self.main_frame = ttk.Frame(self, padding=12)
        self.main_frame.pack(fill="both", expand=True)

        # ---- 顶部工具栏（5 个按钮 + 视频信息开关，排成一行）----
        toolbar = ttk.Frame(self.main_frame)
        toolbar.pack(fill="x", pady=(0, 4))

        self.btn_browse = ttk.Button(toolbar, text="浏览...", bootstyle="primary", command=self.on_browse)
        self.btn_browse.pack(side="left", padx=(0, 6))
        self.buttons.append(self.btn_browse)

        self.btn_refresh = ttk.Button(toolbar, text="刷新", bootstyle="primary", command=self.on_refresh)
        self.btn_refresh.pack(side="left", padx=(0, 12))
        self.buttons.append(self.btn_refresh)

        # 模式切换
        self.mode_var = tk.StringVar(value="模式: 固定")
        self.btn_mode = ttk.Button(toolbar, textvariable=self.mode_var, bootstyle="primary", command=self.on_mode_click)
        self.btn_mode.pack(side="left", padx=(0, 6))
        self.buttons.append(self.btn_mode)

        # 分辨率切换
        self.res_var = tk.StringVar(value="分辨率: 50%")
        self.btn_res = ttk.Button(toolbar, textvariable=self.res_var, bootstyle="primary", command=self.on_res_click)
        self.btn_res.pack(side="left", padx=(0, 6))
        self.buttons.append(self.btn_res)

        # 视频信息叠加开关
        self.show_info_var = tk.BooleanVar(value=True)
        self.btn_show_info = ttk.Checkbutton(
            toolbar, text="视频显示信息", variable=self.show_info_var,
            bootstyle="round-toggle", command=self.on_toggle_show_info,
        )
        self.btn_show_info.pack(side="left", padx=(12, 0))
        self.buttons.append(self.btn_show_info)

        # 文件夹路径信息（单独一行，位于按钮下方，不再挤占工具栏）
        folder_row = ttk.Frame(self.main_frame)
        folder_row.pack(fill="x", pady=(0, 8))
        ttk.Label(folder_row, text="文件夹:").pack(side="left")
        self.folder_var = tk.StringVar(value="未选择")
        ttk.Label(folder_row, textvariable=self.folder_var, foreground="#666666").pack(side="left", padx=(4, 0))

        # ---- 直方图图形对象占位（提前声明，供后续渲染使用）----
        self.hist_bars = None
        self.tf_line = None
        self.low_vline = None
        self.high_vline = None
        self.x_range = None

        # ---- 中部：直方图 + 预览（matplotlib 嵌入，正方形显示） ----
        middle = ttk.Frame(self.main_frame)
        middle.pack(fill="x", pady=(0, 8))  # 不竖向扩展，高度由宽度决定
        middle.columnconfigure(0, weight=1, uniform="cell")
        middle.columnconfigure(1, weight=1, uniform="cell")
        # row 0 不设置 weight！高度由回调锁为 column 的实际宽度
        self._middle_square_job = None
        middle.bind("<Configure>", lambda e: self._schedule_middle_square(middle, e))

        # 直方图（正方形显示区）
        hist_wrap = ttk.LabelFrame(middle, text="直方图 (Histogram)")
        hist_wrap.grid(row=0, column=0, sticky="nsew", padx=(0, 6))
        self.fig_hist = Figure(figsize=(5, 5), dpi=100, facecolor="white")
        self.ax_hist = self.fig_hist.add_subplot(111)
        self.ax_hist.set_axis_off()
        self.fig_hist.subplots_adjust(left=0.02, right=0.98, bottom=0.02, top=0.98)
        self.canvas_hist = FigureCanvasTkAgg(self.fig_hist, master=hist_wrap)
        self.canvas_hist._tkcanvas.place(relwidth=1, relheight=1)  # place 填充，不反向推父容器
        # placeholder 文字用 tkinter Label 叠加（不依赖 matplotlib 中文字体）
        self.hist_placeholder = ttk.Label(hist_wrap, text="请选择文件夹",
                                          font=("", 16), foreground="#cccccc",
                                          anchor="center", background="white")
        self.hist_placeholder.place(relx=0.5, rely=0.5, anchor="center")

        # 预览图（正方形显示区）
        preview_wrap = ttk.LabelFrame(middle, text="预览 (Preview)")
        preview_wrap.grid(row=0, column=1, sticky="nsew", padx=(6, 0))
        self.fig_preview = Figure(figsize=(5, 5), dpi=100, facecolor="white")
        self.ax_preview = self.fig_preview.add_subplot(111)
        self.ax_preview.set_axis_off()
        self.fig_preview.subplots_adjust(left=0.02, right=0.98, bottom=0.02, top=0.98)
        self.canvas_preview = FigureCanvasTkAgg(self.fig_preview, master=preview_wrap)
        self.canvas_preview._tkcanvas.place(relwidth=1, relheight=1)
        self.preview_placeholder = ttk.Label(preview_wrap, text="预览",
                                             font=("", 16), foreground="#cccccc",
                                             anchor="center", background="white")
        self.preview_placeholder.place(relx=0.5, rely=0.5, anchor="center")

        # ---- 中部右侧：滑块（与直方图+预览同一行，右侧竖排） ----
        # 改为把滑块放在 preview 的右边（一个三列 middle）会破坏正方形比例，
        # 所以滑块单独放一行（位于 canvas 下方，不会跑飞）。
        slider_frame = ttk.LabelFrame(self.main_frame, text="亮度 / 对比度调整")
        slider_frame.pack(fill="x", pady=(0, 6))

        self.slider_low_var = tk.DoubleVar(value=1.0)
        self.slider_high_var = tk.DoubleVar(value=99.0)
        self.slider_brightness_var = tk.DoubleVar(value=0.0)
        self.slider_contrast_var = tk.DoubleVar(value=1.0)

        self.low_slider = self._make_slider(slider_frame, "低 Low%", 0, 100, self.slider_low_var, 0.1, "#dd3333")
        self.high_slider = self._make_slider(slider_frame, "高 High%", 0, 100, self.slider_high_var, 0.1, "#33aa33")
        self.brightness_slider = self._make_slider(slider_frame, "亮度", -0.5, 0.5, self.slider_brightness_var, 0.01, "#cc8800")
        self.contrast_slider = self._make_slider(slider_frame, "对比度", 0.1, 5.0, self.slider_contrast_var, 0.01, "#0077cc")

        # ---- 底部按钮行（重置 / 确认并生成PPT / 使用说明 排成一行）----
        btn_row = ttk.Frame(self.main_frame)
        btn_row.pack(fill="x", pady=(6, 0))

        self.btn_reset = ttk.Button(btn_row, text="重置", bootstyle="primary", command=self.on_reset)
        self.btn_reset.pack(side="left", padx=(0, 8))
        self.buttons.append(self.btn_reset)

        self.apply_var = tk.StringVar(value="确认并生成PPT")
        self.btn_apply = ttk.Button(btn_row, textvariable=self.apply_var, bootstyle="primary", command=self.on_apply)
        self.btn_apply.pack(side="left", padx=(0, 8))
        self.buttons.append(self.btn_apply)

        self.add_info_button(btn_row, "使用说明", self.show_usage, "primary")
        self.add_info_button(btn_row, "About", self.show_about, "primary")

        # 打开输出文件夹（放在最右边）
        self.btn_open_folder = ttk.Button(btn_row, text="打开输出文件夹", bootstyle="primary", command=self.on_open_folder)
        self.btn_open_folder.pack(side="right")
        self.buttons.append(self.btn_open_folder)

        # 底部小字
        footer = ttk.Label(self.main_frame, text="群山行云 2026.7.9 21:59",
                           font=("", 10), foreground="#000000", anchor="center")
        footer.pack(side="bottom", fill="x", pady=(8, 0))

        # 启动时主动跑一次 square 约束（Configure 可能未触发）
        middle.after(100, lambda: self._apply_middle_square(middle))

    def _schedule_middle_square(self, middle, _event=None):
        """节流 + 把 row 0 的高度锁成列宽（像素）。"""
        if self._middle_square_job is not None:
            self.after_cancel(self._middle_square_job)
        self._middle_square_job = self.after(20, self._apply_middle_square, middle)

    def _apply_middle_square(self, middle):
        self._middle_square_job = None
        w = middle.winfo_width()
        if w <= 1:
            return
        side = max(150, w // 2 - 12)  # 12 = 两侧 padx 余量
        middle.rowconfigure(0, minsize=side)

    def _make_slider(self, parent, label, vmin, vmax, var, step, color):
        frame = ttk.Frame(parent)
        frame.pack(fill="x", padx=8, pady=2)
        ttk.Label(frame, text=label, width=12).pack(side="left")
        slider = ttk.Scale(frame, from_=vmin, to=vmax, variable=var, command=self._on_slider_change)
        slider.pack(side="left", fill="x", expand=True, padx=(0, 8))
        val_label = ttk.Label(frame, text=f"{var.get():g}", width=8)
        val_label.pack(side="left")
        var._val_label = val_label
        var._color = color
        return slider

    def add_info_button(self, parent, text, command, bootstyle="info-outline"):
        btn = ttk.Button(parent, text=text, command=command, bootstyle=bootstyle)
        btn.pack(side="left", padx=(0, 8))
        self.buttons.append(btn)

    # ---------- 滑块事件 ----------
    def _on_slider_change(self, _val=None):
        now = time.time()
        if now - self.last_update_time < 0.033:
            return
        self.last_update_time = now

        for var in (self.slider_low_var, self.slider_high_var,
                    self.slider_brightness_var, self.slider_contrast_var):
            if hasattr(var, '_val_label'):
                var._val_label.config(text=f"{var.get():g}")

        low_pct = self.slider_low_var.get()
        high_pct = self.slider_high_var.get()
        if low_pct >= high_pct:
            if self.slider_low_var.get() != high_pct - 0.1:
                self.slider_low_var.set(high_pct - 0.1)
            else:
                self.slider_high_var.set(low_pct + 0.1)

        self._refresh_detail()

    def _refresh_detail(self):
        if self.preview_data is None:
            return
        low_pct = self.slider_low_var.get()
        high_pct = self.slider_high_var.get()
        brightness = self.slider_brightness_var.get()
        contrast = self.slider_contrast_var.get()

        sf = self.sorted_flat
        n = len(sf)
        lo = sf[min(n - 1, int(n * max(0.001, low_pct) / 100))]
        hi = sf[min(n - 1, int(n * min(99.999, high_pct) / 100))]

        # 更新 vline
        if self.low_vline is not None:
            self.low_vline.set_xdata([lo, lo])
        if self.high_vline is not None:
            self.high_vline.set_xdata([hi, hi])

        # 更新 tf_line
        if self.tf_line is not None and self.x_range is not None:
            tf_x = np.clip(self.x_range, lo, hi)
            tf_y = (tf_x - lo) / (hi - lo) if hi > lo else tf_x * 0
            tf_y = tf_y + brightness
            tf_y = (tf_y - 0.5) * max(0.01, contrast) + 0.5
            tf_y = np.clip(tf_y, 0, 1)
            self.tf_line.set_ydata(tf_y * self.hist_max_for_scale)

        # 更新预览图
        adjusted, _, _ = apply_bc(self.preview_data, low_pct, high_pct, brightness, contrast)
        self.ax_preview.clear()
        self.ax_preview.imshow(adjusted, cmap="gray", vmin=0, vmax=1, aspect="equal")
        self.ax_preview.set_xticks([])
        self.ax_preview.set_yticks([])
        self.ax_preview.tick_params(length=0)
        self.ax_preview.spines[:].set_color('#cccccc')
        self.ax_preview.spines[:].set_linewidth(0.8)
        self.canvas_preview.draw_idle()
        self.canvas_hist.draw_idle()

    def _hide_placeholders(self):
        """隐藏直方图和预览的 placeholder 文字（文件夹加载后调用）。"""
        if hasattr(self, 'hist_placeholder') and self.hist_placeholder is not None:
            self.hist_placeholder.place_forget()
        if hasattr(self, 'preview_placeholder') and self.preview_placeholder is not None:
            self.preview_placeholder.place_forget()

    # ---------- 直方图刷新 ----------
    def _refresh_preview_display(self, preview_data, slider_init=None):
        """刷新直方图 + 预览图。
        slider_init: 传入 {'low_pct','high_pct','brightness','contrast'} 时使用该初始值
                     (单独模式记忆恢复)；None 则重置为默认值 (1.0/99.0/0.0/1.0)。
        """
        h, w = preview_data.shape
        max_dim = max(h, w)
        if max_dim > 800:
            scale = 800 / max_dim
            preview_data = cv2.resize(preview_data, (int(w * scale), int(h * scale)),
                                      interpolation=cv2.INTER_AREA)

        data_min = preview_data.min()
        data_max = preview_data.max()
        hist_vals, hist_bins = np.histogram(preview_data.flatten(), bins=256, range=(data_min, data_max))
        hist_bin_centers = (hist_bins[:-1] + hist_bins[1:]) / 2
        hist_max = hist_vals.max()

        self.preview_data = preview_data
        self.sorted_flat = np.sort(preview_data.ravel())
        self.hist_max_for_scale = hist_max

        # 加载数据后隐藏 placeholder
        self._hide_placeholders()

        # 决定初始 Low/High：有记忆用记忆，没记忆用 1.0/99.0
        init_lo_pct = slider_init['low_pct'] if slider_init else 1.0
        init_hi_pct = slider_init['high_pct'] if slider_init else 99.0
        init_lo = np.percentile(preview_data, init_lo_pct)
        init_hi = np.percentile(preview_data, init_hi_pct)
        self.x_range = np.linspace(data_min, data_max, 500)

        self.ax_hist.clear()
        self.ax_hist.spines[:].set_color('#cccccc')
        self.ax_hist.spines[:].set_linewidth(0.8)
        self.ax_hist.set_xticks([])
        self.ax_hist.set_yticks([])
        self.ax_hist.tick_params(length=0)
        self.hist_bars = self.ax_hist.bar(hist_bin_centers, hist_vals,
                                           width=(data_max - data_min) / 256,
                                           color='#aaaaaa', alpha=0.8, edgecolor='none')
        tf_x = np.clip(self.x_range, init_lo, init_hi)
        tf_y = (tf_x - init_lo) / (init_hi - init_lo)
        self.tf_line, = self.ax_hist.plot(self.x_range, tf_y * hist_max, '#0077cc', linewidth=2.5)
        self.low_vline = self.ax_hist.axvline(init_lo, color='#dd3333', linestyle='--', linewidth=2)
        self.high_vline = self.ax_hist.axvline(init_hi, color='#33aa33', linestyle='--', linewidth=2)
        self.fig_hist.subplots_adjust(left=0.02, right=0.98, bottom=0.02, top=0.98)
        self.canvas_hist.draw_idle()

        # 重绘预览图（按当前滑条值渲染）
        brightness = slider_init['brightness'] if slider_init else 0.0
        contrast = slider_init['contrast'] if slider_init else 1.0
        adjusted_init, _, _ = apply_bc(preview_data, init_lo_pct, init_hi_pct, brightness, contrast)
        self.ax_preview.clear()
        self.ax_preview.imshow(adjusted_init, cmap='gray', vmin=0, vmax=1, aspect="equal")
        self.ax_preview.set_xticks([])
        self.ax_preview.set_yticks([])
        self.ax_preview.tick_params(length=0)
        self.ax_preview.spines[:].set_color('#cccccc')
        self.ax_preview.spines[:].set_linewidth(0.8)
        self.fig_preview.subplots_adjust(left=0.02, right=0.98, bottom=0.02, top=0.98)
        self.canvas_preview.draw_idle()

        # 滑条复位（使用指定初始值或默认值）
        self.slider_low_var.set(init_lo_pct)
        self.slider_high_var.set(init_hi_pct)
        self.slider_brightness_var.set(brightness)
        self.slider_contrast_var.set(contrast)
        for var in (self.slider_low_var, self.slider_high_var,
                    self.slider_brightness_var, self.slider_contrast_var):
            if hasattr(var, '_val_label'):
                var._val_label.config(text=f"{var.get():g}")

    # ---------- 文件夹加载 ----------
    def load_folder_and_preview(self, folder_path):
        files = [f for f in os.listdir(folder_path) if f.lower().endswith(('.mrc', '.dm4'))]
        if not files:
            messagebox.showwarning("无文件", f"文件夹中没有 mrc/dm4 文件:\n{folder_path}")
            return False

        files = sorted(files, key=lambda x: os.path.getmtime(os.path.join(folder_path, x)))
        earliest_path = os.path.join(folder_path, files[0])

        self.folder = folder_path
        self.files = files
        self.current_file_idx = 0
        self.per_file_params = {}
        self.folder_var.set(folder_path)
        self.title(f"MRC/DM4 → PPT — {os.path.basename(folder_path)}")

        preview_data = load_preview_frame(earliest_path)
        self._refresh_preview_display(preview_data)
        self._update_apply_button()
        return True

    def _load_file_as_preview(self, filepath):
        preview_data = load_preview_frame(filepath)
        # 单独模式：若当前文件已有记忆参数，恢复它；否则用默认初始值
        saved = self.per_file_params.get(os.path.basename(filepath)) if self.mode == 'individual' else None
        params = {
            'low_pct': saved['low_pct'] if saved else 1.0,
            'high_pct': saved['high_pct'] if saved else 99.0,
            'brightness': saved.get('brightness', 0.0) if saved else 0.0,
            'contrast': saved.get('contrast', 1.0) if saved else 1.0,
        }
        self._refresh_preview_display(preview_data, slider_init=params)
        self._update_apply_button()

    # ---------- 工具栏事件 ----------
    def on_browse(self):
        root = tk.Tk()
        root.withdraw()
        root.attributes('-topmost', True)
        folder = filedialog.askdirectory(title='选择包含 MRC/DM4 文件的文件夹')
        root.destroy()
        if folder:
            self.load_folder_and_preview(folder)

    def on_refresh(self):
        if self.folder:
            self.load_folder_and_preview(self.folder)

    def on_mode_click(self):
        self.mode = 'individual' if self.mode == 'fixed' else 'fixed'
        self.mode_var.set(f"模式: {'固定' if self.mode == 'fixed' else '单独'}")
        self._update_apply_button()

    def on_res_click(self):
        options = [25, 50, 75, 100]
        idx = options.index(self.resolution) if self.resolution in options else 1
        self.resolution = options[(idx + 1) % len(options)]
        self.res_var.set(f"分辨率: {self.resolution}%")

    def on_open_folder(self):
        if not self.folder:
            messagebox.showwarning("未选文件夹", "请先选择包含 MRC/DM4 的文件夹。")
            return
        if os.path.isdir(self.folder):
            os.startfile(self.folder)
        else:
            messagebox.showerror("不存在", f"文件夹不存在：{self.folder}")

    def on_toggle_show_info(self):
        """切换视频文件名/帧数信息叠加。"""
        self.show_info = self.show_info_var.get()

    def on_reset(self):
        self.slider_low_var.set(1.0)
        self.slider_high_var.set(99.0)
        self.slider_brightness_var.set(0.0)
        self.slider_contrast_var.set(1.0)
        self._on_slider_change()

    # ---------- 使用说明 ----------
    def show_usage(self):
        messagebox.showinfo(
            "使用说明",
            "1. 点击 [浏览...] 选择包含 *.mrc / *.dm4 的文件夹。\n"
            "2. 调整 Low/High/亮度/对比度滑块观察预览变化。\n"
            "3. 模式: 固定 = 所有文件同一套参数; 单独 = 逐个文件分别调。\n"
            "4. 分辨率 = 输出图片/视频的缩放比，100% 为原图。\n"
            "5. 点击 [确认并生成PPT]，界面会播放Neco-Arc的动画并生成。\n"
            "6. 输出文件保存在同文件夹：mrc_dm4_output_bc.pptx。"
        )

    def show_about(self):
        messagebox.showinfo(
            "About",
            "2026.4.8，开始使用 Python，做的第一个自动化程序。"
            "经过多次小修小补始终达不到完美。\n\n"
            "2026.6.2 引入 Fiji 的直方图，随后陆续添加了视频和 DM4 的支持。\n\n"
            "2026.7.9 将麻烦的 matplotlib 移除，添加记忆功能并美化界面，\n"
            "添加 Neco-Arc 的运行动画，至臻完美。"
        )

    # ---------- 生成参数收集 ----------
    def _current_params(self):
        return {
            'low_pct': self.slider_low_var.get(),
            'high_pct': self.slider_high_var.get(),
            'brightness': self.slider_brightness_var.get(),
            'contrast': self.slider_contrast_var.get(),
        }

    def _update_apply_button(self):
        if self.mode == 'fixed' or not self.files:
            self.apply_var.set("确认并生成PPT")
        else:
            idx = self.current_file_idx
            total = len(self.files)
            if idx >= total - 1:
                self.apply_var.set(f"确认并生成PPT ({total}/{total})")
            else:
                self.apply_var.set(f"保存并继续 ({idx + 1}/{total})")

    # ---------- GIF 动效（hallsep 模板）----------
    def load_gif(self, path):
        frames = []
        if not path.exists():
            return frames
        index = 0
        while True:
            try:
                frame = tk.PhotoImage(file=str(path), format=f"gif -index {index}")
            except tk.TclError:
                break
            frames.append(frame)
            index += 1
        return frames

    def _scale_frame_to(self, frame, target_w, target_h):
        fw = max(frame.width(), 1)
        fh = max(frame.height(), 1)
        if fw <= target_w and fh <= target_h:
            zoom_x = target_w // fw
            zoom_y = target_h // fh
            factor = min(zoom_x, zoom_y)
            if factor > 1:
                return frame.zoom(factor, factor)
            return frame
        else:
            sub_x = max(1, (fw + target_w - 1) // target_w)
            sub_y = max(1, (fh + target_h - 1) // target_h)
            factor = max(sub_x, sub_y)
            if factor > 1:
                return frame.subsample(factor, factor)
            return frame

    def start_gif(self):
        raw_frames = self._preloaded_frames or self.load_gif(GIF_PATH)
        if not raw_frames:
            return
        self.stop_gif()
        if self.main_frame is not None:
            self.main_frame.pack_forget()

        self.run_frame = ttk.Frame(self, padding=14)
        self.run_frame.pack(fill="both", expand=True)
        self.run_frame.columnconfigure(0, weight=1)
        self.run_frame.rowconfigure(0, weight=1)

        self.update_idletasks()
        win_w = max(self.winfo_width(), 1)
        win_h = max(self.winfo_height(), 1)
        target_w = max(100, int(win_w * GIF_DISPLAY_SCALE))
        target_h = max(100, int(win_h * GIF_DISPLAY_SCALE))
        self.run_gif_frames = [self._scale_frame_to(f, target_w, target_h) for f in raw_frames]

        self.run_gif_label = ttk.Label(self.run_frame, image=self.run_gif_frames[0])
        self.run_gif_label.grid(row=0, column=0)

        self.gif_running = True
        self.animate_loop_gif(0)

    def stop_gif(self):
        self.gif_running = False
        if self.gif_after_id is not None:
            self.after_cancel(self.gif_after_id)
            self.gif_after_id = None

    def restore_main_frame(self):
        self.stop_gif()
        if self.run_frame is not None and self.run_frame.winfo_exists():
            self.run_frame.destroy()
        self.run_frame = None
        self.run_gif_label = None
        if self.main_frame is not None and not self.main_frame.winfo_ismapped():
            self.main_frame.pack(fill="both", expand=True)

    def animate_loop_gif(self, index):
        if not self.gif_running or not self.run_gif_frames or self.run_gif_label is None:
            return
        self.run_gif_label.configure(image=self.run_gif_frames[index])
        next_index = (index + 1) % len(self.run_gif_frames)
        self.gif_after_id = self.after(80, self.animate_loop_gif, next_index)

    # ---------- 运行控制 ----------
    def set_buttons_state(self, enabled):
        state = "normal" if enabled else "disabled"
        for btn in self.buttons:
            btn.config(state=state)

    def on_apply(self):
        if self.preview_data is None:
            messagebox.showwarning("未选文件夹", "请先选择包含 MRC/DM4 的文件夹。")
            return

        current = self._current_params()

        if self.mode == 'fixed':
            self.fixed_params = current
            self._kickoff_run()
        else:
            cur_file = self.files[self.current_file_idx]
            self.per_file_params[cur_file] = current
            print(f"已保存: {cur_file}  → Low={current['low_pct']:.1f}% High={current['high_pct']:.1f}%")

            next_idx = self.current_file_idx + 1
            if next_idx >= len(self.files):
                self.fixed_params = current
                self._kickoff_run()
            else:
                self.current_file_idx = next_idx
                next_path = os.path.join(self.folder, self.files[next_idx])
                self._load_file_as_preview(next_path)

    def _kickoff_run(self):
        self.set_buttons_state(False)
        self.start_gif()
        threading.Thread(target=self._run_ppt, daemon=True).start()

    def _progress_cb(self, done, total, name):
        self.after(0, self.apply_var.set, f"生成中 {done}/{total}")

    def _run_ppt(self):
        try:
            out_path = generate_ppt(
                folder=self.folder,
                files=self.files,
                mode=self.mode,
                resolution=self.resolution,
                per_file_params=self.per_file_params,
                default_params=self.fixed_params,
                progress_cb=self._progress_cb,
                show_info=self.show_info,
            )
            self.after(0, self._on_ppt_success, out_path)
        except Exception as exc:
            self.after(0, self._on_ppt_error, exc)

    def _on_ppt_success(self, out_path):
        self.play_done_once()
        self.set_buttons_state(True)
        self.apply_var.set("确认并生成PPT")
        messagebox.showinfo("完成", f"PPT 已生成：\n{out_path}")

    def _on_ppt_error(self, exc):
        self.restore_main_frame()
        self.set_buttons_state(True)
        self.apply_var.set("确认并生成PPT")
        messagebox.showerror("运行失败", str(exc))

    def play_done_once(self):
        self.stop_gif()
        self.after(1200, self.restore_main_frame)

    def run(self):
        self.mainloop()


if __name__ == "__main__":
    App().run()