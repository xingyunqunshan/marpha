import os
import sys
from tkinter import Tk
from tkinter.filedialog import askdirectory

import discretisedfield as df
import matplotlib.pyplot as plt
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Cm

IMAGE_TRANSPOSE = getattr(Image, "Transpose", Image)

SUPPORTED_EXTENSIONS = (".omf", ".ovf")


def is_supported_field_file(name):
    return name.lower().endswith(SUPPORTED_EXTENSIONS)


def has_field_file(folder):
    return any(is_supported_field_file(name) for name in os.listdir(folder))


def find_field_folders(parent):
    if not os.path.isdir(parent):
        raise FileNotFoundError(f"Folder does not exist: {parent}")

    folders = []
    for root, _, _ in os.walk(parent):
        if has_field_file(root):
            folders.append(root)

    return sorted(folders, key=os.path.getmtime)


def process_folder(folder):
    files = [f for f in os.listdir(folder) if is_supported_field_file(f)]
    files = sorted(files, key=lambda x: os.path.getmtime(os.path.join(folder, x)))

    if not files:
        print(f"Skip, no OMF/OVF files: {folder}")
        return

    print(f"\nProcessing folder: {folder}")
    print(f"Found {len(files)} OMF/OVF files")

    second_file_path = os.path.join(folder, files[1]) if len(files) >= 2 else folder

    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    cols = 6
    img_h = Cm(5.5)
    img_w = Cm(5.6)

    x_gap = Cm(0.01)
    left_margin = Cm(0.1)
    top_margin = Cm(1.5)

    for idx, file in enumerate(files):

        if idx % 12 == 0:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            textbox_top = slide.shapes.add_textbox(Cm(1), Cm(0.2), Cm(30), Cm(0.6))
            textbox_top.text = second_file_path

        filepath = os.path.join(folder, file)

        print(f"[{idx + 1}/{len(files)}] Reading: {filepath}")

        # -- load OMF / OVF via discretisedfield -------------------------------
        # OVF comes in two flavours (text v1.0, binary/text v2.0); the library
        # handles both, but older installs may only expose the deprecated API.
        try:
            field = df.Field.from_file(filepath)
        except AttributeError:
            try:
                field = df.Field.fromfile(filepath)  # deprecated fallback
            except AttributeError:
                print(f"  ERROR: cannot read {filepath} – discretisedfield "
                      f"is too old or the file format is unsupported.")
                continue
        except Exception as exc:
            print(f"  ERROR: failed to parse '{filepath}': {exc}")
            continue

        data = field.array

        print("Data shape:", data.shape)

        if data.ndim != 4 or data.shape[-1] != 3:
            raise ValueError(f"Unexpected OMF/OVF data shape: {data.shape}")

        # =====================================================
        # 中层切片
        # =====================================================
        z = data.shape[2] // 2
        mx = data[:, :, z, 0]
        my = data[:, :, z, 1]
        mz = data[:, :, z, 2]

        # =====================================================
        # Mz 背景归一化（保持你原逻辑）
        # =====================================================
        low = np.percentile(mz, 1)
        high = np.percentile(mz, 99)

        if high > low:
            mz_show = np.clip(mz, low, high)
            mz_show = (mz_show - low) / (high - low)
        else:
            mz_show = mz

        # =====================================================
        # ⭐关键修复：磁矩归一化（解决全黑）
        # =====================================================
        norm = np.sqrt(mx**2 + my**2 + mz**2)
        norm[norm == 0] = 1

        mx_n = mx / norm
        my_n = my / norm

        # =====================================================
        # 箭头网格
        # =====================================================
        step = max(1, mx.shape[0] // 100)

        x_grid, y_grid = np.meshgrid(
            np.arange(0, mx.shape[1], step),
            np.arange(0, mx.shape[0], step),
        )

        u = my_n[::step, ::step]
        v = mx_n[::step, ::step]

        # =====================================================
        # 画图（箭头风格修复版）
        # =====================================================
        img_path = os.path.join(folder, f"temp_{idx}.png")

        plt.figure(figsize=(6, 6), dpi=400)

        plt.imshow(mz_show, cmap="bwr", origin="lower")

        plt.quiver(
            x_grid,
            y_grid,
            u,
            v,
            color="black",
            scale=120,
            width=0.002,
            headwidth=2,
            headlength=2,
            headaxislength=2,
            pivot="mid"
        )

        plt.axis("off")

        plt.savefig(img_path, bbox_inches="tight", pad_inches=0)
        plt.close()

        # =====================================================
        # 图像翻转（保持你原来的PPT视觉）
        # =====================================================
        with Image.open(img_path) as img:
            img = img.transpose(IMAGE_TRANSPOSE.FLIP_LEFT_RIGHT)
            img = img.transpose(IMAGE_TRANSPOSE.ROTATE_270)
            img.save(img_path)

        # =====================================================
        # PPT 排版（完全不动）
        # =====================================================
        pos = idx % 12
        row = pos // cols
        col = pos % cols

        x = left_margin + col * (img_w + x_gap)
        y = top_margin if row == 0 else top_margin + Cm(9)
        title_y = y if col % 2 == 0 else y + Cm(1.2)

        name = os.path.splitext(file)[0]

        textbox = slide.shapes.add_textbox(x, title_y, img_w, Cm(0.5))
        textbox.text = name

        slide.shapes.add_picture(
            img_path,
            x,
            y + Cm(2),
            height=img_h,
        )

    out_path = os.path.join(folder, "omf_ovf_top_surface_output.pptx")
    prs.save(out_path)

    print(f"Saved: {out_path}")


def main():
    # -- pop up a folder picker ------------------------------------------------
    root = Tk()
    root.withdraw()          # hide the root window
    root.attributes("-topmost", True)  # bring dialog to front

    parent = askdirectory(
        title="选择包含 OMF/OVF 文件的文件夹",
    )
    root.destroy()

    if not parent:
        print("未选择文件夹，已取消。")
        return 0

    print(f"Scanning folder: {parent}")
    folders = find_field_folders(parent)

    if not folders:
        print("No OMF/OVF files found.")
        return 1

    for folder in folders:
        process_folder(folder)

    print("\nAll done.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
