#将文件夹内部的所有PPT做成一个大的ppt
import os
import tkinter as tk
from tkinter import filedialog
from copy import deepcopy
from pptx import Presentation

# ======================
# 手动选择文件夹
# ======================
root_dir = tk.Tk()
root_dir.withdraw()  # 隐藏主窗口
root = filedialog.askdirectory(title="选择包含PPT的文件夹")
root_dir.destroy()

if not root:
    print("未选择文件夹，程序退出。")
    exit()

# ======================
# 搜索ppt
# ======================
ppt_files = []

for folderpath, _, filenames in os.walk(root):
    for f in filenames:
        if f.lower().endswith(".pptx") and not f.startswith("~$"):
            ppt_files.append(os.path.join(folderpath, f))

# ======================
# 按时间排序
# ======================
ppt_files.sort(key=lambda x: os.path.getmtime(x))

# ======================
# 新PPT
# ======================
first = Presentation(ppt_files[0])

merged = Presentation()
merged.slide_width = first.slide_width
merged.slide_height = first.slide_height

# 删除默认页
if len(merged.slides) > 0:
    rId = merged.slides._sldIdLst[0].rId
    merged.part.drop_rel(rId)
    del merged.slides._sldIdLst[0]

# ======================
# 复制 slide（含图片）
# ======================
def copy_slide(source_slide, dest_prs):

    blank_layout = dest_prs.slide_layouts[6]
    new_slide = dest_prs.slides.add_slide(blank_layout)

    for shape in source_slide.shapes:

        # 图片
        if shape.shape_type == 13:
            image = shape.image
            image_bytes = image.blob

            temp_path = "temp_img.png"
            with open(temp_path, "wb") as f:
                f.write(image_bytes)

            new_slide.shapes.add_picture(
                temp_path,
                shape.left,
                shape.top,
                shape.width,
                shape.height
            )

            os.remove(temp_path)

        else:
            el = deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')

# ======================
# 合并
# ======================
for i, ppt_path in enumerate(ppt_files):

    print("处理:", ppt_path)

    prs = Presentation(ppt_path)

    for slide in prs.slides:
        copy_slide(slide, merged)

    # 插空白页
    if i < len(ppt_files) - 1:
        merged.slides.add_slide(merged.slide_layouts[6])

# ======================
# 保存
# ======================
out_path = os.path.join(root, "汇总PPT.pptx")
merged.save(out_path)

print("完成：", out_path)