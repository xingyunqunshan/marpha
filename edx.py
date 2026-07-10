import zipfile
import io
import queue
import threading
import traceback
from datetime import datetime
from pathlib import Path

try:
    from lxml import etree
except ImportError as exc:
    etree = None
    LXML_IMPORT_ERROR = exc
else:
    LXML_IMPORT_ERROR = None

try:
    from pptx import Presentation
    from pptx.util import Cm, Inches, Pt
except ImportError as exc:
    Presentation = None
    Cm = Inches = Pt = None
    PPTX_IMPORT_ERROR = exc
else:
    PPTX_IMPORT_ERROR = None

try:
    import tkinter as tk
    from tkinter import filedialog, messagebox
except ImportError as exc:
    tk = None
    filedialog = None
    messagebox = None
    TK_IMPORT_ERROR = exc
else:
    TK_IMPORT_ERROR = None

try:
    import ttkbootstrap as ttk
except ImportError as exc:
    if tk is not None:
        from tkinter import ttk
    else:
        ttk = None
    TTKBOOTSTRAP_IMPORT_ERROR = exc
else:
    TTKBOOTSTRAP_IMPORT_ERROR = None


# =========================================================
# Word 文件
# =========================================================
DOCX_PATH = r"C:\Users\理塘丁真\Desktop\新建文件夹 (3)\Project 1_Site 21_2026-05-20_14-37-38PTTE2 - 副本.docx"

# 输出PPT
OUTPUT_PPT = r"C:\Users\理塘丁真\Desktop\output.pptx"


# =========================================================
# XML 命名空间
# =========================================================
NS = {
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
}


def get_cell_text(cell):

    return "".join(
        cell.xpath(
            ".//w:t/text()",
            namespaces=NS
        )
    ).strip()


def parse_word_table(table):

    rows = []

    for tr in table.xpath("./w:tr", namespaces=NS):

        row = []

        for cell in tr.xpath("./w:tc", namespaces=NS):
            row.append(get_cell_text(cell))

        rows.append(row)

    return rows


def parse_spectrum_table(sdt, fallback_texts):

    tables = sdt.xpath(
        ".//w:tbl",
        namespaces=NS
    )

    if not tables:
        return None

    rows = parse_word_table(tables[0])

    title = fallback_texts[0]
    header_index = None

    for idx, row in enumerate(rows):

        if row and row[0].startswith("Spectrum"):
            title = row[0]

        if row and row[0] == "Element":
            header_index = idx
            break

    if header_index is None:
        return None

    headers = rows[header_index]
    col_count = len(headers)
    body_rows = []

    for row in rows[header_index + 1:]:

        fixed_row = row[:col_count]

        if len(fixed_row) < col_count:
            fixed_row.extend([""] * (col_count - len(fixed_row)))

        if any(cell != "" for cell in fixed_row):
            body_rows.append(fixed_row)

    return {
        "type": "TABLE",
        "title": title,
        "headers": headers,
        "rows": body_rows,
        "formula": ""
    }


def parse_float(value):

    try:
        return float(str(value).replace("%", "").strip())
    except ValueError:
        return None


def get_formula_items(table_data):

    headers = table_data["headers"]

    if "Element" not in headers or "Atomic %" not in headers:
        return []

    element_index = headers.index("Element")
    atomic_index = headers.index("Atomic %")

    items = []

    for row in table_data["rows"]:

        if element_index >= len(row) or atomic_index >= len(row):
            continue

        element = row[element_index].strip()

        if not element or element == "Total":
            continue

        atomic_percent = parse_float(row[atomic_index])

        if atomic_percent is None:
            continue

        items.append((element, atomic_percent))

    return items


def collect_formula_elements(stream):

    elements = []

    for item in stream:

        if item["type"] != "TABLE":
            continue

        for element, _ in get_formula_items(item):

            if element not in elements:
                elements.append(element)

    return elements


def choose_docx_path():

    try:
        import tkinter as tk
        from tkinter import filedialog, messagebox
    except Exception:
        return DOCX_PATH

    root = tk.Tk()
    root.withdraw()

    path = filedialog.askopenfilename(
        title="选择 Word 文件",
        filetypes=[
            ("Word 文件", "*.docx"),
            ("所有文件", "*.*")
        ]
    )

    root.destroy()

    if not path:
        try:
            messagebox.showinfo(
                "已取消",
                "没有选择 Word 文件，程序结束。"
            )
        except Exception:
            pass

        return None

    return path


def make_output_path(docx_path):

    path = Path(docx_path)

    return str(
        path.with_name(f"{path.stem}_output.pptx")
    )


def choose_formula_options(elements):

    if not elements:
        return None

    try:
        import tkinter as tk
        from tkinter import messagebox, ttk
    except Exception:
        return {
            "basis_element": elements[0],
            "basis_value": 1.0
        }

    result = {
        "basis_element": elements[0],
        "basis_value": 1.0
    }

    root = tk.Tk()
    root.title("选择化学式基准")
    root.resizable(False, False)

    tk.Label(root, text="固定哪个元素").grid(
        row=0,
        column=0,
        padx=12,
        pady=(12, 6),
        sticky="w"
    )

    element_var = tk.StringVar(value=elements[0])
    element_box = ttk.Combobox(
        root,
        textvariable=element_var,
        values=elements,
        state="readonly",
        width=18
    )
    element_box.grid(
        row=0,
        column=1,
        padx=12,
        pady=(12, 6)
    )

    tk.Label(root, text="固定数值").grid(
        row=1,
        column=0,
        padx=12,
        pady=6,
        sticky="w"
    )

    value_var = tk.StringVar(value="1")
    value_entry = ttk.Entry(
        root,
        textvariable=value_var,
        width=20
    )
    value_entry.grid(
        row=1,
        column=1,
        padx=12,
        pady=6,
        sticky="w"
    )

    def confirm():

        try:
            basis_value = float(value_var.get())
        except ValueError:
            messagebox.showerror(
                "输入错误",
                "固定数值必须是数字，例如 1、2、1.5。"
            )
            return

        if basis_value <= 0:
            messagebox.showerror(
                "输入错误",
                "固定数值必须大于 0。"
            )
            return

        result["basis_element"] = element_var.get()
        result["basis_value"] = basis_value
        root.destroy()

    ttk.Button(
        root,
        text="确定",
        command=confirm
    ).grid(
        row=2,
        column=0,
        columnspan=2,
        padx=12,
        pady=(8, 12)
    )

    root.mainloop()

    return result


def format_formula_number(value):

    text = f"{value:.3f}".rstrip("0").rstrip(".")

    return text if text else "0"


def build_formula(table_data, formula_options):

    if formula_options is None:
        return ""

    items = get_formula_items(table_data)

    if not items:
        return ""

    basis_element = formula_options["basis_element"]
    basis_value = formula_options["basis_value"]
    basis_atomic = None

    for element, atomic_percent in items:

        if element == basis_element:
            basis_atomic = atomic_percent
            break

    if not basis_atomic:
        return ""

    parts = []

    for element, atomic_percent in items:

        amount = atomic_percent / basis_atomic * basis_value
        parts.append(
            f"{element}{format_formula_number(amount)}"
        )

    return "".join(parts)


def build_formula_amounts(table_data, formula_options):

    if formula_options is None:
        return {}

    items = get_formula_items(table_data)

    if not items:
        return {}

    basis_element = formula_options["basis_element"]
    basis_value = formula_options["basis_value"]
    basis_atomic = None

    for element, atomic_percent in items:

        if element == basis_element:
            basis_atomic = atomic_percent
            break

    if not basis_atomic:
        return {}

    amounts = {}

    for element, atomic_percent in items:
        amounts[element] = atomic_percent / basis_atomic * basis_value

    return amounts


def apply_formula_options(stream, formula_options):

    for item in stream:

        if item["type"] == "TABLE":
            item["formula"] = build_formula(
                item,
                formula_options
            )
            item["formula_amounts"] = build_formula_amounts(
                item,
                formula_options
            )


# =========================================================
# 解析 Word
# 目标：
# IMAGE
# TEXT
# IMAGE
# TEXT
# ...
#
# 严格按 XML 顺序
# 只保留 PNG
# =========================================================
def parse_word(docx_path):

    import zipfile
    from lxml import etree

    NS = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    }

    # =====================================================
    # 元素周期表
    # 用于自动识别表格开始
    # =====================================================
    ELEMENTS = {
        "H","He","Li","Be","B","C","N","O","F","Ne",
        "Na","Mg","Al","Si","P","S","Cl","Ar","K","Ca",
        "Sc","Ti","V","Cr","Mn","Fe","Co","Ni","Cu","Zn",
        "Ga","Ge","As","Se","Br","Kr","Rb","Sr","Y","Zr",
        "Nb","Mo","Tc","Ru","Rh","Pd","Ag","Cd","In","Sn",
        "Sb","Te","I","Xe","Cs","Ba","La","Ce","Pr","Nd",
        "Pm","Sm","Eu","Gd","Tb","Dy","Ho","Er","Tm","Yb",
        "Lu","Hf","Ta","W","Re","Os","Ir","Pt","Au","Hg",
        "Tl","Pb","Bi","Po","At","Rn"
    }

    stream = []

    with zipfile.ZipFile(docx_path) as z:

        # =================================================
        # document.xml
        # =================================================
        xml_data = z.read("word/document.xml")
        root = etree.fromstring(xml_data)

        # =================================================
        # rels
        # =================================================
        rel_data = z.read("word/_rels/document.xml.rels")
        rel_root = etree.fromstring(rel_data)

        # rId -> media path
        rid_map = {}

        for rel in rel_root:

            rid = rel.get("Id")
            target = rel.get("Target")

            rid_map[rid] = target

        # =================================================
        # body
        # =================================================
        body = root.xpath("//w:body", namespaces=NS)[0]

        # =================================================
        # 遍历 body
        # =================================================
        for child in body:

            # =================================================
            # IMAGE
            # =================================================
            if child.tag.endswith("p"):

                drawings = child.xpath(
                    ".//w:drawing",
                    namespaces=NS
                )

                if drawings:

                    for d in drawings:

                        blip = d.xpath(
                            ".//a:blip",
                            namespaces=NS
                        )

                        if not blip:
                            continue

                        rId = blip[0].get(
                            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                        )

                        if rId not in rid_map:
                            continue

                        media_path = "word/" + rid_map[rId]

                        # 只保留 PNG
                        if not media_path.lower().endswith(".png"):
                            continue

                        try:
                            img_bytes = z.read(media_path)
                        except:
                            continue

                        stream.append({
                            "type": "IMAGE",
                            "bytes": img_bytes
                        })

            # =================================================
            # SDT（EDX表）
            # =================================================
            elif child.tag.endswith("sdt"):

                texts = child.xpath(
                    ".//w:t/text()",
                    namespaces=NS
                )

                texts = [
                    t.strip()
                    for t in texts
                    if t.strip()
                ]

                if not texts:
                    continue

                # =============================================
                # Spectrum 表
                # =============================================
                if texts[0].startswith("Spectrum"):

                    table_data = parse_spectrum_table(
                        child,
                        texts
                    )

                    if table_data is not None:
                        stream.append(table_data)
                    else:
                        stream.append({
                            "type": "TEXT",
                            "text": "\n".join(texts)
                        })

                # =============================================
                # 普通 SDT
                # =============================================
                else:

                    stream.append({
                        "type": "TEXT",
                        "text": "\n".join(texts)
                    })

    return stream
# =========================================================
# 分页逻辑
#
# 每 2 张图 -> 一页 PPT
# 后续连续 TEXT -> 当前页
# =========================================================
def make_chunks(stream):

    chunks = []

    current_images = []
    current_text = []
    current_tables = []

    for item in stream:

        # =============================================
        # IMAGE
        # =============================================
        if item["type"] == "IMAGE":

            # 已经2图
            # 开新页
            if len(current_images) == 2:

                chunks.append({
                    "images": current_images,
                    "text": current_text,
                    "tables": current_tables
                })

                current_images = []
                current_text = []
                current_tables = []

            current_images.append(item["bytes"])

        # =============================================
        # TEXT
        # =============================================
        elif item["type"] == "TEXT":

            current_text.append(item["text"])

        # =============================================
        # TABLE
        # =============================================
        elif item["type"] == "TABLE":

            current_tables.append(item)

    # 最后一页
    if current_images or current_text or current_tables:

        chunks.append({
            "images": current_images,
            "text": current_text,
            "tables": current_tables
        })

    return chunks


def set_cell_text(cell, text, font_size=7, bold=False):

    cell.text = str(text)
    paragraph = cell.text_frame.paragraphs[0]

    if paragraph.runs:
        run = paragraph.runs[0]
    else:
        run = paragraph.add_run()

    run.font.size = Pt(font_size)
    run.font.bold = bold


def add_formula_label(slide, formula):

    if not formula:
        return

    textbox = slide.shapes.add_textbox(
        Inches(0.3),
        Inches(0.02),
        Inches(4.4),
        Inches(0.28)
    )

    textbox.text_frame.text = formula
    run = textbox.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(18)
    run.font.bold = True


def add_real_table(slide, table_data, left, top, width, max_height):

    headers = table_data["headers"]
    rows = table_data["rows"]

    if not headers:
        return 0

    title_box = slide.shapes.add_textbox(
        left,
        top,
        width,
        Inches(0.25)
    )
    title_box.text_frame.text = table_data["title"]
    title_run = title_box.text_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(10)
    title_run.font.bold = True

    row_count = len(rows) + 1
    col_count = len(headers)

    table_top = top + Inches(0.3)
    row_height = Inches(0.22)
    table_height = min(
        row_height * row_count,
        max_height - Inches(0.3)
    )

    shape = slide.shapes.add_table(
        row_count,
        col_count,
        left,
        table_top,
        width,
        table_height
    )

    table = shape.table

    col_width = int(width / col_count)

    for col in table.columns:
        col.width = col_width

    for col_idx, header in enumerate(headers):
        set_cell_text(
            table.cell(0, col_idx),
            header,
            font_size=7,
            bold=True
        )

    for row_idx, row in enumerate(rows, start=1):

        for col_idx in range(col_count):

            value = row[col_idx] if col_idx < len(row) else ""

            set_cell_text(
                table.cell(row_idx, col_idx),
                value,
                font_size=7
            )

    return Inches(0.35) + table_height


def save_presentation(prs, output_path):

    try:
        prs.save(output_path)
        return output_path
    except PermissionError:
        path = Path(output_path)
        fallback_path = path.with_name(
            f"{path.stem}_{datetime.now():%Y%m%d_%H%M%S}{path.suffix}"
        )
        prs.save(fallback_path)
        return str(fallback_path)


def collect_tables_from_chunks(chunks):

    tables = []

    for chunk in chunks:
        tables.extend(chunk.get("tables", []))

    return tables


def build_summary(elements, tables):

    rows = []

    for table_data in tables:

        amounts = table_data.get("formula_amounts", {})

        if not amounts:
            continue

        rows.append({
            "name": table_data.get("title", ""),
            "amounts": amounts
        })

    averages = {}

    for element in elements:

        values = [
            row["amounts"][element]
            for row in rows
            if element in row["amounts"]
        ]

        if values:
            averages[element] = sum(values) / len(values)
        else:
            averages[element] = None

    return rows, averages


def add_summary_table(slide, elements, rows, averages):

    table_rows = len(rows) + 2
    table_cols = len(elements) + 1

    shape = slide.shapes.add_table(
        table_rows,
        table_cols,
        Inches(0.4),
        Inches(0.45),
        Inches(9.2),
        Inches(0.32 * table_rows)
    )

    table = shape.table
    table.columns[0].width = Inches(1.35)

    data_col_width = int(Inches(7.85) / len(elements))

    for col_idx in range(1, table_cols):
        table.columns[col_idx].width = data_col_width

    set_cell_text(
        table.cell(0, 0),
        "Spectrum",
        font_size=8,
        bold=True
    )

    for col_idx, element in enumerate(elements, start=1):
        set_cell_text(
            table.cell(0, col_idx),
            element,
            font_size=8,
            bold=True
        )

    for row_idx, row in enumerate(rows, start=1):

        set_cell_text(
            table.cell(row_idx, 0),
            row["name"],
            font_size=8
        )

        for col_idx, element in enumerate(elements, start=1):

            value = row["amounts"].get(element)

            set_cell_text(
                table.cell(row_idx, col_idx),
                "" if value is None else format_formula_number(value),
                font_size=8
            )

    average_row = len(rows) + 1

    set_cell_text(
        table.cell(average_row, 0),
        "Average",
        font_size=8,
        bold=True
    )

    for col_idx, element in enumerate(elements, start=1):

        value = averages.get(element)

        set_cell_text(
            table.cell(average_row, col_idx),
            "" if value is None else format_formula_number(value),
            font_size=8,
            bold=True
        )


def add_summary_slides(prs, chunks):

    tables = collect_tables_from_chunks(chunks)

    elements = []

    for table_data in tables:

        for element in table_data.get("formula_amounts", {}):

            if element not in elements:
                elements.append(element)

    if not elements:
        return

    rows, averages = build_summary(elements, tables)

    if not rows:
        return

    max_elements_per_slide = 8
    max_rows_per_slide = 16

    for element_start in range(0, len(elements), max_elements_per_slide):

        element_group = elements[
            element_start:element_start + max_elements_per_slide
        ]

        group_averages = {
            element: averages.get(element)
            for element in element_group
        }

        for row_start in range(0, len(rows), max_rows_per_slide):

            row_group = rows[
                row_start:row_start + max_rows_per_slide
            ]

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            title_box = slide.shapes.add_textbox(
                Inches(0.4),
                Inches(0.1),
                Inches(9.2),
                Inches(0.3)
            )
            title_box.text_frame.text = "Element Ratio Summary"
            title_run = title_box.text_frame.paragraphs[0].runs[0]
            title_run.font.size = Pt(14)
            title_run.font.bold = True

            add_summary_table(
                slide,
                element_group,
                row_group,
                group_averages
            )


def collect_second_images(chunks):

    second_images = []

    for idx, chunk in enumerate(chunks, start=1):

        images = chunk.get("images", [])

        if len(images) < 2:
            continue

        label = f"Slide {idx}"

        if chunk.get("tables"):
            label = chunk["tables"][0].get("title", label)

        second_images.append({
            "label": label,
            "bytes": images[1]
        })

    return second_images


def add_picture_fit(slide, image_bytes, left, top, max_width, max_height):

    picture = slide.shapes.add_picture(
        io.BytesIO(image_bytes),
        left,
        top,
        width=max_width
    )

    if picture.height > max_height:

        ratio = max_height / picture.height
        picture.height = max_height
        picture.width = int(picture.width * ratio)

    picture.left = left + int((max_width - picture.width) / 2)
    picture.top = top + int((max_height - picture.height) / 2)


def add_second_image_slides(prs, chunks):

    second_images = collect_second_images(chunks)

    if not second_images:
        return

    per_slide = 12
    cols = 4
    rows = 3

    margin_left = Inches(0.25)
    margin_top = Inches(0.55)
    cell_width = Inches(2.35)
    cell_height = Inches(2.05)
    gap_x = Inches(0.12)
    gap_y = Inches(0.12)
    label_height = Inches(0.18)

    for start in range(0, len(second_images), per_slide):

        group = second_images[start:start + per_slide]
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(
            Inches(0.35),
            Inches(0.12),
            Inches(9.2),
            Inches(0.3)
        )
        title_box.text_frame.text = "Second Image Summary"
        title_run = title_box.text_frame.paragraphs[0].runs[0]
        title_run.font.size = Pt(14)
        title_run.font.bold = True

        for item_index, item in enumerate(group):

            row = item_index // cols
            col = item_index % cols

            left = margin_left + col * (cell_width + gap_x)
            top = margin_top + row * (cell_height + gap_y)

            add_picture_fit(
                slide,
                item["bytes"],
                left,
                top,
                cell_width,
                cell_height - label_height
            )

            label_box = slide.shapes.add_textbox(
                left,
                top + cell_height - label_height,
                cell_width,
                label_height
            )
            label_box.text_frame.text = item["label"]
            label_run = label_box.text_frame.paragraphs[0].runs[0]
            label_run.font.size = Pt(7)


# =========================================================
# 生成 PPT
# =========================================================
def build_ppt(chunks, output_path):

    prs = Presentation()

    add_summary_slides(prs, chunks)
    add_second_image_slides(prs, chunks)

    for idx, chunk in enumerate(chunks):

        slide = prs.slides.add_slide(
            prs.slide_layouts[6]
        )

        imgs = chunk["images"]

        if chunk.get("tables"):
            add_formula_label(
                slide,
                chunk["tables"][0].get("formula", "")
            )

        # =================================================
        # 左图
        # =================================================
        if len(imgs) >= 1:

            slide.shapes.add_picture(
                io.BytesIO(imgs[0]),
                Inches(0.3),
                Inches(0.3),
                width=Inches(4.5)
            )

        # =================================================
        # 右图
        # =================================================
        if len(imgs) >= 2:

            slide.shapes.add_picture(
                io.BytesIO(imgs[1]),
                Inches(5),
                Inches(0.3),
                width=Inches(4.5)
            )

        # =================================================
        # 下方文字
        # =================================================
        content_top = Inches(3.5)

        if chunk["text"]:

            text = "\n".join(chunk["text"])

            textbox = slide.shapes.add_textbox(
                Inches(0.5),
                content_top,
                Inches(9),
                Inches(0.8)
            )

            textbox.text_frame.text = text
            content_top += Inches(0.9)

        # =================================================
        # 真表格
        # =================================================
        content_top += Cm(2)

        for table_data in chunk.get("tables", []):

            used_height = add_real_table(
                slide,
                table_data,
                Inches(0.5),
                content_top,
                Inches(9),
                Inches(6.9) - content_top
            )

            content_top += used_height + Inches(0.15)

        print(f"完成 slide {idx+1}")

    return save_presentation(prs, output_path)


# =========================================================
# GUI MAIN
# =========================================================
def get_missing_dependencies():

    missing = []

    if LXML_IMPORT_ERROR is not None:
        missing.append(("lxml", LXML_IMPORT_ERROR))

    if PPTX_IMPORT_ERROR is not None:
        missing.append(("python-pptx", PPTX_IMPORT_ERROR))

    if TK_IMPORT_ERROR is not None:
        missing.append(("tkinter", TK_IMPORT_ERROR))

    if TTKBOOTSTRAP_IMPORT_ERROR is not None:
        missing.append(("ttkbootstrap", TTKBOOTSTRAP_IMPORT_ERROR))

    return missing


def summarize_stream(stream):

    return {
        "stream": len(stream),
        "images": sum(1 for item in stream if item["type"] == "IMAGE"),
        "text": sum(1 for item in stream if item["type"] == "TEXT"),
        "tables": sum(1 for item in stream if item["type"] == "TABLE"),
        "elements": collect_formula_elements(stream),
    }


def run_edx_conversion(docx_path, output_path, formula_options, log=print):

    log(f"Word 文件: {docx_path}")
    log(f"输出 PPT: {output_path}")
    log("解析 Word...")

    stream = parse_word(docx_path)
    summary = summarize_stream(stream)

    log(f"Stream 数量: {summary['stream']}")
    log(f"PNG 图片数量: {summary['images']}")
    log(f"TEXT 数量: {summary['text']}")
    log(f"TABLE 数量: {summary['tables']}")

    if summary["elements"]:
        log("可选公式元素: " + ", ".join(summary["elements"]))

    apply_formula_options(stream, formula_options)

    log("分页处理中...")
    chunks = make_chunks(stream)
    log(f"PPT 页数: {len(chunks)}")

    log("生成 PPT...")
    actual_output = build_ppt(chunks, output_path)
    log("完成")
    log(f"输出: {actual_output}")

    return actual_output


def get_window_base():

    if TTKBOOTSTRAP_IMPORT_ERROR is None:
        return ttk.Window

    if tk is not None:
        return tk.Tk

    return object


class EdxApp(get_window_base()):

    def __init__(self):

        if TTKBOOTSTRAP_IMPORT_ERROR is None:
            super().__init__(
                title="EDX 数据转 PPT",
                themename="flatly",
                size=(860, 620),
                resizable=(True, True)
            )
        else:
            super().__init__()
            self.title("EDX 数据转 PPT")
            self.geometry("860x620")
            self.resizable(True, True)

        self.docx_var = tk.StringVar()
        self.output_var = tk.StringVar()
        self.basis_element_var = tk.StringVar()
        self.basis_value_var = tk.StringVar(value="1")
        self.status_var = tk.StringVar(value="请选择 Word 文件")
        self.stream = None
        self.elements = []
        self.message_queue = queue.Queue()
        self.worker = None

        self.build_ui()
        self.after(100, self.drain_queue)

    def button_options(self, bootstyle=None):

        if TTKBOOTSTRAP_IMPORT_ERROR is None and bootstyle:
            return {"bootstyle": bootstyle}

        return {}

    def build_ui(self):

        root = ttk.Frame(self, padding=18)
        root.pack(fill="both", expand=True)
        root.columnconfigure(0, weight=1)
        root.rowconfigure(4, weight=1)

        header = ttk.Frame(root)
        header.grid(row=0, column=0, sticky="ew", pady=(0, 14))
        header.columnconfigure(0, weight=1)

        ttk.Label(
            header,
            text="EDX 数据转 PPT",
            font=("Microsoft YaHei UI", 18, "bold")
        ).grid(row=0, column=0, sticky="w")

        ttk.Label(
            header,
            textvariable=self.status_var,
            foreground="#5c6670"
        ).grid(row=1, column=0, sticky="w", pady=(4, 0))

        self.add_path_row(
            root,
            1,
            "Word 文件",
            self.docx_var,
            self.choose_docx,
            "选择"
        )
        self.add_path_row(
            root,
            2,
            "输出 PPT",
            self.output_var,
            self.choose_output,
            "另存为"
        )

        options = ttk.Labelframe(root, text="化学式基准", padding=12)
        options.grid(row=3, column=0, sticky="ew", pady=(12, 12))
        options.columnconfigure(1, weight=1)

        ttk.Label(options, text="固定元素").grid(row=0, column=0, sticky="w")
        self.element_box = ttk.Combobox(
            options,
            textvariable=self.basis_element_var,
            state="readonly",
            width=18
        )
        self.element_box.grid(row=0, column=1, sticky="w", padx=(10, 24))

        ttk.Label(options, text="固定数值").grid(row=0, column=2, sticky="w")
        ttk.Entry(
            options,
            textvariable=self.basis_value_var,
            width=12
        ).grid(row=0, column=3, sticky="w", padx=(10, 0))

        actions = ttk.Frame(root)
        actions.grid(row=5, column=0, sticky="ew", pady=(12, 0))
        actions.columnconfigure(3, weight=1)

        self.analyze_button = ttk.Button(
            actions,
            text="分析 Word",
            command=self.start_analyze,
            **self.button_options("info")
        )
        self.analyze_button.grid(row=0, column=0, padx=(0, 8))

        self.generate_button = ttk.Button(
            actions,
            text="生成 PPT",
            command=self.start_generate,
            **self.button_options("success")
        )
        self.generate_button.grid(row=0, column=1, padx=(0, 8))

        self.open_output_button = ttk.Button(
            actions,
            text="打开输出目录",
            command=self.open_output_folder,
            state="disabled"
        )
        self.open_output_button.grid(row=0, column=2, padx=(0, 8))

        self.log_text = tk.Text(
            root,
            height=14,
            wrap="word",
            relief="solid",
            borderwidth=1,
            bg="#ffffff",
            fg="#1f2933"
        )
        self.log_text.grid(row=4, column=0, sticky="nsew")

    def add_path_row(self, parent, row, label, variable, command, button_text):

        frame = ttk.Frame(parent)
        frame.grid(row=row, column=0, sticky="ew", pady=4)
        frame.columnconfigure(1, weight=1)

        ttk.Label(frame, text=label, width=10).grid(row=0, column=0, sticky="w")
        ttk.Entry(frame, textvariable=variable).grid(
            row=0,
            column=1,
            sticky="ew",
            padx=(10, 8)
        )
        ttk.Button(frame, text=button_text, command=command).grid(
            row=0,
            column=2
        )

    def choose_docx(self):

        path = filedialog.askopenfilename(
            title="选择 Word 文件",
            filetypes=[("Word 文件", "*.docx"), ("所有文件", "*.*")]
        )

        if not path:
            return

        self.docx_var.set(path)
        self.output_var.set(make_output_path(path))
        self.stream = None
        self.elements = []
        self.update_elements([])
        self.open_output_button.config(state="disabled")
        self.status_var.set("已选择 Word 文件，请先分析")

    def choose_output(self):

        initial = self.output_var.get()
        initial_dir = str(Path(initial).parent) if initial else ""
        initial_file = Path(initial).name if initial else "output.pptx"

        path = filedialog.asksaveasfilename(
            title="保存 PPT",
            defaultextension=".pptx",
            initialdir=initial_dir,
            initialfile=initial_file,
            filetypes=[("PowerPoint 文件", "*.pptx"), ("所有文件", "*.*")]
        )

        if path:
            self.output_var.set(path)

    def validate_docx_path(self):

        path = self.docx_var.get().strip()

        if not path:
            raise ValueError("请先选择 Word 文件。")

        if not Path(path).exists():
            raise ValueError(f"Word 文件不存在: {path}")

        if Path(path).suffix.lower() != ".docx":
            raise ValueError("请选择 .docx 文件。")

        if not self.output_var.get().strip():
            self.output_var.set(make_output_path(path))

        return path

    def validate_formula_options(self):

        if not self.elements:
            return None

        element = self.basis_element_var.get().strip()

        if element not in self.elements:
            raise ValueError("请选择有效的固定元素。")

        try:
            value = float(self.basis_value_var.get())
        except ValueError as exc:
            raise ValueError("固定数值必须是数字，例如 1、2、0.5。") from exc

        if value <= 0:
            raise ValueError("固定数值必须大于 0。")

        return {
            "basis_element": element,
            "basis_value": value
        }

    def start_analyze(self):

        try:
            docx_path = self.validate_docx_path()
        except ValueError as exc:
            messagebox.showerror("参数错误", str(exc))
            return

        self.run_worker(self.analyze_worker, docx_path)

    def analyze_worker(self, docx_path):

        self.queue_status("正在解析 Word...")
        self.queue_log(f"Word 文件: {docx_path}")
        stream = parse_word(docx_path)
        summary = summarize_stream(stream)
        self.message_queue.put(("analysis_done", stream, summary))

    def start_generate(self):

        try:
            docx_path = self.validate_docx_path()
            if self.stream is None:
                raise ValueError("请先点击“分析 Word”，确认元素后再生成 PPT。")
            output_path = self.output_var.get().strip()
            formula_options = self.validate_formula_options()
        except ValueError as exc:
            messagebox.showerror("参数错误", str(exc))
            return

        self.run_worker(
            self.generate_worker,
            docx_path,
            output_path,
            formula_options
        )

    def generate_worker(self, docx_path, output_path, formula_options):

        actual_output = run_edx_conversion(
            docx_path,
            output_path,
            formula_options,
            log=self.queue_log
        )
        self.message_queue.put(("generate_done", actual_output))

    def run_worker(self, target, *args):

        if self.worker and self.worker.is_alive():
            messagebox.showinfo("正在运行", "当前任务还没有结束，请稍后。")
            return

        self.set_busy(True)
        self.log_text.delete("1.0", "end")

        self.worker = threading.Thread(
            target=self.worker_wrapper,
            args=(target, args),
            daemon=True
        )
        self.worker.start()

    def worker_wrapper(self, target, args):

        try:
            target(*args)
        except Exception as exc:
            self.message_queue.put(("error", exc, traceback.format_exc()))
        finally:
            self.message_queue.put(("idle",))

    def queue_log(self, text):

        self.message_queue.put(("log", text))

    def queue_status(self, text):

        self.message_queue.put(("status", text))

    def drain_queue(self):

        while True:
            try:
                message = self.message_queue.get_nowait()
            except queue.Empty:
                break

            kind = message[0]

            if kind == "log":
                self.append_log(message[1])
            elif kind == "status":
                self.status_var.set(message[1])
            elif kind == "analysis_done":
                self.handle_analysis_done(message[1], message[2])
            elif kind == "generate_done":
                self.handle_generate_done(message[1])
            elif kind == "error":
                self.handle_error(message[1], message[2])
            elif kind == "idle":
                self.set_busy(False)

        self.after(100, self.drain_queue)

    def append_log(self, text):

        self.log_text.insert("end", f"{text}\n")
        self.log_text.see("end")

    def handle_analysis_done(self, stream, summary):

        self.stream = stream
        self.elements = summary["elements"]
        self.update_elements(self.elements)

        self.append_log(f"Stream 数量: {summary['stream']}")
        self.append_log(f"PNG 图片数量: {summary['images']}")
        self.append_log(f"TEXT 数量: {summary['text']}")
        self.append_log(f"TABLE 数量: {summary['tables']}")

        if self.elements:
            self.append_log("可选公式元素: " + ", ".join(self.elements))
            self.status_var.set("分析完成，请确认化学式基准后生成 PPT")
        else:
            self.append_log("未识别到可计算化学式的元素，将不生成公式标签。")
            self.status_var.set("分析完成，可直接生成 PPT")

    def update_elements(self, elements):

        self.element_box.config(values=elements)

        if elements:
            self.basis_element_var.set(elements[0])
        else:
            self.basis_element_var.set("")

    def handle_generate_done(self, actual_output):

        self.output_var.set(actual_output)
        self.open_output_button.config(state="normal")
        self.status_var.set("PPT 生成完成")
        messagebox.showinfo("完成", f"PPT 已生成:\n{actual_output}")

    def handle_error(self, exc, detail):

        self.append_log(detail)
        self.status_var.set("运行失败")
        messagebox.showerror("运行失败", str(exc))

    def set_busy(self, busy):

        state = "disabled" if busy else "normal"
        self.analyze_button.config(state=state)
        self.generate_button.config(state=state)

    def open_output_folder(self):

        output = self.output_var.get().strip()

        if not output:
            return

        folder = Path(output).parent

        if not folder.exists():
            messagebox.showerror("目录不存在", str(folder))
            return

        import os

        os.startfile(folder)


def show_dependency_error(missing):

    details = "\n".join(
        f"- {name}: {exc}"
        for name, exc in missing
    )

    message = (
        "EDX 转 PPT 缺少必要依赖，程序无法运行。\n\n"
        f"{details}\n\n"
        "请确认当前 Python 环境已安装: lxml、python-pptx、ttkbootstrap。"
    )

    if tk is not None and messagebox is not None:
        root = tk.Tk()
        root.withdraw()
        messagebox.showerror("依赖缺失", message)
        root.destroy()
    else:
        print(message)


def main():

    missing = get_missing_dependencies()

    if missing:
        show_dependency_error(missing)
        raise SystemExit(1)

    EdxApp().mainloop()


if __name__ == "__main__":
    main()
